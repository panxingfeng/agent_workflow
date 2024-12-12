import asyncio
import json
import os
import re
from typing import List, Dict, Optional
from dataclasses import dataclass
from enum import Enum
import logging

import whisper
from lightrag import QueryParam
from lightrag.lightrag import LightRAG
from lightrag.llm import ollama_model_complete, ollama_embedding

from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
from bs4 import BeautifulSoup
from lightrag.utils import EmbeddingFunc
import chardet

from agent_workflow.rag.base import BaseRAG
from config.config import OLLAMA_DATA


def get_project_root() -> str:
    current_path = os.path.abspath(__file__)
    project_root = os.path.dirname(os.path.dirname(os.path.dirname(current_path)))
    logging.info(f"Project root path: {project_root}")
    return project_root


def get_upload_dir(path_name:str = None) -> str:
    """获取上传文件夹路径"""
    if path_name is None:
        path = "upload"
    else:
        path = os.path.join(path_name,"upload")
    upload_dir = os.path.join(get_project_root(), path)
    if not os.path.exists(upload_dir):
        os.makedirs(upload_dir)
    return upload_dir


def get_files_from_upload(path_name:str = None) -> List[str]:
    upload_dir = get_upload_dir(path_name)
    logging.info(f"扫描上传目录: {upload_dir}")
    supported_files = []
    supported_extensions = {ft.value for ft in FileType}

    for filename in os.listdir(upload_dir):
        ext = os.path.splitext(filename)[1].lower().lstrip('.')
        filepath = os.path.join(upload_dir, filename)
        logging.info(f"检查文件: {filepath}")

        if os.path.isfile(filepath) and ext in supported_extensions:
            supported_files.append(filepath)
            logging.info(f"找到支持的文件: {filepath}")

    return sorted(supported_files)


class FileType(Enum):
    """支持的文件类型枚举"""
    DOC = 'doc'
    DOCX = 'docx'
    PDF = 'pdf'
    PPT = 'ppt'
    PPTX = 'pptx'
    TXT = 'txt'
    HTML = 'html'
    HTM = 'htm'
    MP3 = 'mp3'
    WAV = 'wav'
    MD = 'md'
    JSON = 'json'

@dataclass
class ProcessResult:
    """文件处理结果"""
    filename: str
    success: bool
    content: Optional[str] = None
    error: Optional[str] = None
    file_type: Optional[FileType] = None


class DocumentProcessor(BaseRAG):
    """文档处理和RAG知识库管理器"""
    def __init__(self, path_name: str = "document_rag",files_path_name:str = None):
        self.output_dir = os.path.join(get_project_root(), path_name)
        self.rag = None
        self.logger = logging.getLogger(__name__)
        self.whisper_model = None
        self.files_path_name = files_path_name

    def _setup_rag(self):
        """初始化RAG实例"""
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir)

        self.rag = LightRAG(
            working_dir=self.output_dir,
            llm_model_func=ollama_model_complete,
            llm_model_name=OLLAMA_DATA['model'],
            llm_model_kwargs={"options": {"num_ctx": 32768}},
            embedding_func=EmbeddingFunc(
                embedding_dim=1024,
                max_token_size=8192,
                func=lambda texts: ollama_embedding(texts, embed_model=OLLAMA_DATA['embedding_model'])
            ),
        )

    def _load_whisper_model(self):
        """加载Whisper模型"""
        if self.whisper_model is None:
            self.whisper_model = whisper.load_model("turbo")

    def _get_file_type(self, file_path: str) -> Optional[FileType]:
        """获取文件类型"""
        ext = os.path.splitext(file_path)[1].lower().lstrip('.')
        try:
            return FileType(ext)
        except ValueError:
            return None

    def _process_docx(self, file_path: str) -> str:
        try:
            logging.info(f"Attempting to open docx file: {file_path}")
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")
            doc = Document(file_path)
            return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            logging.error(f"Error processing docx file {file_path}: {str(e)}")
            raise

    def _process_pdf(self, file_path: str) -> str:
        """处理PDF文件"""
        reader = PdfReader(file_path)
        text = []
        for page in reader.pages:
            text.append(page.extract_text())
        return '\n'.join(text)

    def _process_pptx(self, file_path: str) -> str:
        """处理PPTX文件"""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)

    def _process_txt(self, file_path: str) -> str:
        """处理TXT文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError as e:
            self.logger.error(f"文件 {file_path} 编码错误: {str(e)}")
            raise

    def _process_html(self, file_path: str) -> str:
        """处理HTML文件"""
        with open(file_path, 'rb') as file:
            raw_data = file.read()
            result = chardet.detect(raw_data)
            encoding = result['encoding']

        with open(file_path, 'r', encoding=encoding) as file:
            soup = BeautifulSoup(file, 'html.parser')
            # 移除script和style元素
            for script in soup(["script", "style"]):
                script.decompose()
            # 获取文本
            text = soup.get_text()
            # 处理空白字符
            lines = (line.strip() for line in text.splitlines())
            return '\n'.join(line for line in lines if line)

    def _process_audio(self, file_path: str) -> str:
        """处理音频文件"""
        self._load_whisper_model()
        result = self.whisper_model.transcribe(file_path)
        return result["text"]

    def _process_md(self, file_path: str) -> str:
        """处理 Markdown 文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
                # 移除 Markdown 语法标记
                # 移除代码块
                content = re.sub(r'```[\s\S]*?```', '', content)
                # 移除行内代码
                content = re.sub(r'`[^`]*`', '', content)
                # 移除标题标记
                content = re.sub(r'#{1,6}\s', '', content)
                # 移除链接，保留文本
                content = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', content)
                # 移除图片
                content = re.sub(r'!\[([^\]]*)\]\([^\)]+\)', '', content)
                # 移除粗体和斜体
                content = re.sub(r'\*\*([^\*]*)\*\*', r'\1', content)
                content = re.sub(r'\*([^\*]*)\*', r'\1', content)
                content = re.sub(r'__([^_]*)__', r'\1', content)
                content = re.sub(r'_([^_]*)_', r'\1', content)
                # 移除列表标记
                content = re.sub(r'^\s*[-*+]\s', '', content, flags=re.MULTILINE)
                content = re.sub(r'^\s*\d+\.\s', '', content, flags=re.MULTILINE)

                return content.strip()
        except UnicodeDecodeError as e:
            self.logger.error(f"文件 {file_path} 编码错误: {str(e)}")
            raise

    def _process_json(self, file_path: str) -> str:
        """处理 JSON 文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = json.load(file)
                # 将 JSON 内容转换为字符串表示
                if isinstance(content, (dict, list)):
                    # 美化输出 JSON 内容
                    return json.dumps(content, ensure_ascii=False, indent=2)
                else:
                    return str(content)
        except json.JSONDecodeError as e:
            self.logger.error(f"JSON 解析错误 {file_path}: {str(e)}")
            raise
        except UnicodeDecodeError as e:
            self.logger.error(f"文件编码错误 {file_path}: {str(e)}")
            raise

    async def _process_file(self, file_path: str) -> ProcessResult:
        """异步处理单个文件"""
        global file_type
        if not os.path.exists(file_path):
            self.logger.error(f"文件不存在: {file_path}")
            return ProcessResult(
                filename=os.path.basename(file_path),
                success=False,
                error=f"文件不存在: {file_path}"
            )

        try:
            file_type = self._get_file_type(file_path)
            if not file_type:
                raise ValueError(f"不支持的文件类型: {os.path.splitext(file_path)[1]}")

            # 根据文件类型选择处理方法
            content = None
            if file_type == FileType.DOCX:
                content = await asyncio.to_thread(self._process_docx, file_path)
            elif file_type == FileType.PDF:
                content = await asyncio.to_thread(self._process_pdf, file_path)
            elif file_type == FileType.PPTX:
                content = await asyncio.to_thread(self._process_pptx, file_path)
            elif file_type == FileType.MD:
                content = await asyncio.to_thread(self._process_md, file_path)
            elif file_type in [FileType.TXT]:
                content = await asyncio.to_thread(self._process_txt, file_path)
            elif file_type in [FileType.HTML, FileType.HTM]:
                content = await asyncio.to_thread(self._process_html, file_path)
            elif file_type in [FileType.MP3, FileType.WAV]:
                content = await asyncio.to_thread(self._process_audio, file_path)
            elif file_type == FileType.JSON:
                content = await asyncio.to_thread(self._process_json, file_path)
            else:
                raise ValueError(f"暂不支持处理 {file_type.value} 类型的文件")

            if content:
                return ProcessResult(
                    filename=os.path.basename(file_path),
                    success=True,
                    content=content,
                    file_type=file_type
                )

        except Exception as e:
            self.logger.error(f"处理文件 {os.path.basename(file_path)} 时发生错误: {str(e)}")
            return ProcessResult(
                filename=os.path.basename(file_path),
                success=False,
                error=str(e),
                file_type=file_type if 'file_type' in locals() else None
            )

    async def process_documents_async(self, input_files: List[str]) -> Dict[str, List[ProcessResult]]:
        """异步批量处理文档"""
        self._setup_rag()
        results = {'success': [], 'failed': []}
        processed_contents = []

        # 并发处理所有文件
        tasks = [self._process_file(file_path) for file_path in input_files]
        file_results = await asyncio.gather(*tasks)

        # 整理处理结果
        for result in file_results:
            if result.success:
                results['success'].append(result)
                processed_contents.append(result.content)
            else:
                results['failed'].append(result)

        # 批量插入知识库
        if processed_contents:
            try:
                await asyncio.to_thread(self.rag.insert, processed_contents)
                self.logger.info(f"成功将 {len(processed_contents)} 个文档插入知识库")
            except Exception as e:
                self.logger.error(f"插入知识库时发生错误: {str(e)}")
                for result in results['success']:
                    result.success = False
                    result.error = f"插入知识库失败: {str(e)}"
                results['failed'].extend(results.pop('success'))

        return results

    async def run(self):
        """启动处理流程"""
        try:
            input_files = get_files_from_upload(self.files_path_name)

            if not input_files:
                self.logger.warning("upload文件夹中没有找到支持的文件类型")
                self.logger.info(f"支持的文件类型: {', '.join(ft.value for ft in FileType)}")
                return

            # 显示待处理文件
            supported_files = "\n".join(
                f"{os.path.basename(f)} ({os.path.splitext(f)[1].lower().lstrip('.')})"
                for f in input_files
            )
            self.logger.info(f"找到以下支持的文件：\n{supported_files}")

            # 异步处理文件
            results = await self.process_documents_async(input_files)

            # 打印处理摘要
            success_count = len(results['success'])
            failed_count = len(results['failed'])
            total_count = success_count + failed_count

            print(f"\n处理结果:")
            if results['success']:
                print("\n成功处理的文件:")
                for result in results['success']:
                    print(f"- {result.filename}")

            if results['failed']:
                print("\n处理失败的文件:")
                for result in results['failed']:
                    print(
                        f"- {result.filename} ({result.file_type.value if result.file_type else '未知类型'}): {result.error}")

            print(f"\n统计信息:")
            print(f"总文件数: {total_count}")
            print(f"成功数量: {success_count}")
            print(f"失败数量: {failed_count}")
            print(f"成功率: {(success_count / total_count * 100 if total_count else 0):.2f}%")

        except Exception as e:
            self.logger.error(f"处理过程中发生错误: {str(e)}", exc_info=True)
            raise

    def cleanup(self):
        """清理所有资源"""
        try:
            # 清理 whisper 模型
            if hasattr(self, 'whisper_model') and self.whisper_model is not None:
                del self.whisper_model
                self.whisper_model = None

            # 调用父类清理方法
            super().cleanup()

        except Exception as e:
            self.logger.error(f"清理DocumentProcessor资源时出错: {str(e)}")
            raise


class LightsRAG(BaseRAG):
    """RAG问答类"""
    def __init__(self, path_name: str = "document_rag"):
        """
        初始化问答系统

        Args:
            rag_path: RAG知识库路径
        """
        self.output_dir = os.path.join(get_project_root(), path_name)
        self.rag = None
        self.logger = logging.getLogger(__name__)
        self._setup_rag()

    def _setup_rag(self):
        """初始化RAG实例"""
        if not os.path.exists(self.output_dir):
            raise ValueError(f"RAG路径不存在: {self.output_dir}")

        self.rag = LightRAG(
            working_dir=self.output_dir,
            llm_model_func=ollama_model_complete,
            llm_model_name=OLLAMA_DATA['model'],
            llm_model_kwargs={"options": {"num_ctx": 32768}},
            embedding_func=EmbeddingFunc(
                embedding_dim=1024,
                max_token_size=8192,
                func=lambda texts: ollama_embedding(texts, embed_model=OLLAMA_DATA['embedding_model'])
            ),
        )

    async def ask(self, question: str, mode: str = "global") -> str:
        """
        异步方式查询问题答案

        Args:
            question: 用户问题
            mode: 查询模式，默认为"global"

        Returns:
            str: 回答内容
        """
        try:
            answer = await asyncio.to_thread(
                self.rag.query,
                question,
                param=QueryParam(mode=mode)
            )
            return answer
        except Exception as e:
            self.logger.error(f"RAG查询失败: {str(e)}")
            raise
