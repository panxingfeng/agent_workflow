# -*- coding: utf-8 -*-
"""
@author: [PanXingFeng]
@contact: [1115005803@qq.comã€canomiguelittle@gmail.com]
@date: 2025-1-11
@version: 2.0.0
@license: MIT License
Copyright (c) 2024 [PanXingFeng]
All rights reserved.
"""
import asyncio
import json
import os
import platform
import traceback
from enum import Enum
from typing import List, Optional
import tempfile
import zipfile
import subprocess
import sys
import requests
from pptx import Presentation
from pdf2image import convert_from_path
from pptx.util import Inches
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

from agent_workflow.tools.tool.base import BaseTool
from agent_workflow.utils import loadingInfo

# è®¾ç½®Ghostscriptç¯å¢ƒå˜é‡
os.environ["PATH"] += r";C:\Program Files\gs\gs10.04.0\bin"

logger = loadingInfo("file_converter_tool")

class ConversionType(str, Enum):
    """
    æ–‡ä»¶è½¬æ¢ç±»å‹æšä¸¾

    å®šä¹‰æ”¯æŒçš„æ‰€æœ‰è½¬æ¢ç±»å‹ï¼š
    - URL_TO_PDF: ç½‘é¡µè½¬PDF
    - PDF_TO_WORD: PDFè½¬Word
    - PDF_TO_TEXT: PDFè½¬æ–‡æœ¬
    - PDF_TO_IMAGE: PDFè½¬å›¾ç‰‡
    - PDF_TO_PPT: PDFè½¬PPT
    - PDF_TO_MARKDOWN: PDFè½¬Markdown
    - FILE_TO_PDF: å…¶ä»–æ ¼å¼è½¬PDF
    - MARKDOWN_TO_PDF: Markdownè½¬PDF
    """
    URL_TO_PDF = "url_to_pdf"
    PDF_TO_WORD = "pdf_to_word"
    PDF_TO_TEXT = "pdf_to_text"
    PDF_TO_IMAGE = "pdf_to_image"
    PDF_TO_PPT = "pdf_to_ppt"
    PDF_TO_MARKDOWN = "pdf_to_markdown"
    FILE_TO_PDF = "file_to_pdf"
    MARKDOWN_TO_PDF = "markdown_to_pdf"

    @classmethod
    def list_tasks(cls) -> List[str]:
        """
        åˆ—å‡ºæ‰€æœ‰æ”¯æŒçš„è½¬æ¢ç±»å‹

        Returns:
            æ‰€æœ‰æ”¯æŒçš„è½¬æ¢ç±»å‹åˆ—è¡¨
        """
        return [task.value for task in cls]


def outputData(data: str, printInfo: bool = False):
    """
    è¾“å‡ºæ•°æ®çš„è¾…åŠ©å‡½æ•°

    Args:
        data: è¦è¾“å‡ºçš„æ•°æ®
        printInfo: æ˜¯å¦æ‰“å°ä¿¡æ¯
    """
    if printInfo:
        print(data)


class FileConverterTool(BaseTool):
    """
    æ–‡ä»¶è½¬æ¢å·¥å…·ç±»

    åŠŸèƒ½ï¼š
    1. æ”¯æŒå¤šç§æ–‡ä»¶æ ¼å¼çš„è½¬æ¢
    2. è‡ªåŠ¨ç®¡ç†ä¾èµ–å’Œèµ„æº
    3. æä¾›ç»Ÿä¸€çš„è½¬æ¢æ¥å£
    4. æ”¯æŒè‡ªå®šä¹‰è¾“å‡ºæ ¼å¼

    å±æ€§ï¼š
        output_directory: è¾“å‡ºç›®å½•
        printInfo: æ˜¯å¦æ‰“å°è¯¦ç»†ä¿¡æ¯
        base_dir: åŸºç¡€ç›®å½•è·¯å¾„
        poppler_path: Popplerå·¥å…·è·¯å¾„
    """

    def __init__(self, output_directory: str = "output", printInfo: bool = False):
        """
        åˆå§‹åŒ–æ–‡ä»¶è½¬æ¢å·¥å…·

        Args:
            output_directory: è¾“å‡ºç›®å½•è·¯å¾„
            printInfo: æ˜¯å¦æ‰“å°è¯¦ç»†ä¿¡æ¯
        """
        self.upload_dir = "upload"
        self.sub_upload_dir = os.path.join(self.upload_dir, "files")
        self.output_dir = "output"
        self.output_directory = output_directory
        self.printInfo = printInfo
        self.base_dir = os.path.dirname(os.path.abspath(__file__))
        self.poppler_path = os.path.join(self.base_dir, "poppler", "bin")

        # åˆå§‹åŒ–ç›®å½•
        self._ensure_directories()

        # æ£€æŸ¥Poppler
        if not self._check_poppler():
            self._download_poppler()

        self.register_fonts()


    def get_description(self) -> str:
        tool_info = {
            "name": "FileConverterTool",
            "description": "æ–‡ä»¶æ ¼å¼è½¬æ¢å·¥å…·ï¼Œæ”¯æŒå¤šç§æ ¼å¼äº’è½¬",
            "parameters": {
                "conversion_type": {
                    "type": "string",
                    "description": "è½¬æ¢ç±»å‹",
                    "required": True,
                    "enum": [
                        {"name": "url_to_pdf", "description": "ç½‘é¡µè½¬PDF"},
                        {"name": "pdf_to_word", "description": "PDFè½¬Word"},
                        {"name": "pdf_to_text", "description": "PDFè½¬æ–‡æœ¬"},
                        {"name": "pdf_to_image", "description": "PDFè½¬å›¾ç‰‡"},
                        {"name": "pdf_to_ppt", "description": "PDFè½¬PPT"},
                        {"name": "pdf_to_markdown", "description": "PDFè½¬Markdown"},
                        {"name": "file_to_pdf", "description": "å…¶ä»–æ ¼å¼è½¬PDF"},
                        {"name": "markdown_to_pdf", "description": "Markdownè½¬PDF"}
                    ]
                },
                "input_path": {
                    "type": "string",
                    "description": "è¾“å…¥æ–‡ä»¶è·¯å¾„æˆ–URL",
                    "required": True
                }
            }
        }
        return json.dumps(tool_info, ensure_ascii=False)


    def register_fonts(self):
        """
        æ³¨å†Œå­—ä½“

        åŠŸèƒ½ï¼š
        1. å°è¯•æ³¨å†Œç³»ç»Ÿä¸­çš„ä¸­æ–‡å­—ä½“
        2. æ”¯æŒå¤šå¹³å°ï¼ˆWindows/Linux/macOSï¼‰
        3. ä¼˜é›…é™çº§å¤„ç†
        """
        try:
            # å¸¸ç”¨å­—ä½“è·¯å¾„åˆ—è¡¨
            font_paths = [
                # Windowså­—ä½“
                "C:/Windows/Fonts/msyh.ttf",  # å¾®è½¯é›…é»‘
                "C:/Windows/Fonts/simsun.ttc",  # å®‹ä½“
                # Linuxå­—ä½“
                "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf",
                # macOSå­—ä½“
                "/System/Library/Fonts/PingFang.ttc"
            ]

            # å°è¯•æ³¨å†Œç¬¬ä¸€ä¸ªå¯ç”¨çš„å­—ä½“
            for font_path in font_paths:
                if os.path.exists(font_path):
                    try:
                        pdfmetrics.registerFont(TTFont('CustomFont', font_path))
                        return
                    except:
                        continue
        except:
            pass


    def _ensure_directories(self):
        """
        åˆ›å»ºå¿…è¦çš„ç›®å½•ç»“æ„

        åŠŸèƒ½ï¼š
        1. ç¡®ä¿è¾“å‡ºç›®å½•å­˜åœ¨
        2. åˆ›å»ºPopplerå·¥å…·ç›®å½•
        3. æä¾›æ“ä½œåé¦ˆ
        """
        directories = [
            self.output_directory,
            os.path.join(self.base_dir, "poppler")
        ]
        for directory in directories:
            if not os.path.exists(directory):
                os.makedirs(directory)
                logger.info(f"Created directory: {directory}")


    def _check_poppler(self) -> bool:
        """
        æ£€æŸ¥Poppleræ˜¯å¦å·²æ­£ç¡®å®‰è£…

        Returns:
            bool: Poppleræ˜¯å¦å¯ç”¨

        æ£€æŸ¥é¡¹ç›®ï¼š
        1. æ ¸å¿ƒæ‰§è¡Œæ–‡ä»¶å­˜åœ¨æ€§
        2. æ–‡ä»¶æƒé™
        3. ç‰ˆæœ¬å…¼å®¹æ€§
        """
        required_files = ['pdfinfo.exe', 'pdftocairo.exe']
        return all(os.path.exists(os.path.join(self.poppler_path, file))
                   for file in required_files)


    def _download_poppler(self):
        """
        ä¸‹è½½å¹¶å®‰è£…Poppler

        åŠŸèƒ½ï¼š
        1. ä»GitHubä¸‹è½½æœ€æ–°ç‰ˆæœ¬
        2. è‡ªåŠ¨è§£å‹å’Œé…ç½®
        3. é…ç½®ç¯å¢ƒå˜é‡
        4. æä¾›è¯¦ç»†çš„å®‰è£…åé¦ˆ

        å¼‚å¸¸å¤„ç†ï¼š
        - ä¸‹è½½å¤±è´¥æä¾›å¤‡é€‰æ–¹æ¡ˆ
        - è§£å‹é”™è¯¯å¤„ç†
        - æƒé™é—®é¢˜å¤„ç†
        """
        logger.info("\nDownloading Poppler...")

        try:
            # ä¸‹è½½Poppler
            poppler_url = "https://github.com/oschwartz10612/poppler-windows/releases/download/v23.11.0-0/Release-23.11.0-0.zip"
            response = requests.get(poppler_url, stream=True)
            response.raise_for_status()

            # ä¿å­˜å’Œè§£å‹
            temp_zip = os.path.join(self.base_dir, "poppler_temp.zip")
            with open(temp_zip, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)

            logger.info("Extracting Poppler...")
            with zipfile.ZipFile(temp_zip, 'r') as zip_ref:
                poppler_dir = os.path.join(self.base_dir, "poppler")
                zip_ref.extractall(poppler_dir)

            # æ¸…ç†ä¸´æ—¶æ–‡ä»¶
            os.remove(temp_zip)
            logger.info("Poppler installation completed!")

            # ç¯å¢ƒå˜é‡é…ç½®
            if sys.platform == 'win32':
                os.environ['PATH'] = f"{self.poppler_path};{os.environ['PATH']}"

        except Exception as e:
            logger.info(f"Failed to download Poppler: {str(e)}")
            logger.info("Please download manually from: https://github.com/oschwartz10612/poppler-windows/releases")
            logger.info(f"And extract to: {self.poppler_path}")


    def _generate_output_path(self, prefix: str, extension: str) -> str:
        """ç”Ÿæˆè¾“å‡ºæ–‡ä»¶è·¯å¾„"""
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{prefix}_{timestamp}.{extension}"
        return os.path.join(self.output_dir, filename)


    def convert_with_libreoffice(self, input_file: str, output_file: str) -> bool:
        """
        ä½¿ç”¨LibreOfficeè¿›è¡Œæ–‡ä»¶è½¬æ¢

        Args:
            input_file: è¾“å…¥æ–‡ä»¶è·¯å¾„
            output_file: è¾“å‡ºæ–‡ä»¶è·¯å¾„

        Returns:
            bool: è½¬æ¢æ˜¯å¦æˆåŠŸ

        åŠŸèƒ½ï¼š
        1. è‡ªåŠ¨æ£€æµ‹LibreOfficeå®‰è£…
        2. æ”¯æŒå¤šå¹³å°
        3. å‘½ä»¤è¡Œå‚æ•°ä¼˜åŒ–
        4. é”™è¯¯å¤„ç†å’Œé‡è¯•
        """
        try:
            # æ£€æµ‹LibreOfficeè·¯å¾„
            if platform.system() == 'Windows':
                soffice_path = 'C:\\Program Files\\LibreOffice\\program\\soffice.exe'
                if not os.path.exists(soffice_path):
                    soffice_path = 'C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe'
            else:
                soffice_path = 'soffice'

            # æ„å»ºè½¬æ¢å‘½ä»¤
            cmd = [
                soffice_path,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', os.path.dirname(output_file),
                input_file
            ]

            # æ‰§è¡Œè½¬æ¢
            process = subprocess.Popen(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE
            )
            stdout, stderr = process.communicate()

            # æ£€æŸ¥è¾“å‡ºæ–‡ä»¶
            expected_output = os.path.join(
                os.path.dirname(output_file),
                os.path.splitext(os.path.basename(input_file))[0] + '.pdf'
            )

            if os.path.exists(expected_output):
                # é‡å‘½ååˆ°ç›®æ ‡è·¯å¾„
                os.rename(expected_output, output_file)
                return True
            return False

        except Exception as e:
            logger.error(f"LibreOffice conversion failed: {str(e)}")
            return False


    def convert_with_unoconv(self, input_file: str, output_file: str) -> bool:
        """
        ä½¿ç”¨Unoconvè¿›è¡Œæ–‡ä»¶è½¬æ¢

        Args:
            input_file: è¾“å…¥æ–‡ä»¶è·¯å¾„
            output_file: è¾“å‡ºæ–‡ä»¶è·¯å¾„

        Returns:
            bool: è½¬æ¢æ˜¯å¦æˆåŠŸ
        """
        try:
            cmd = ['unoconv', '-f', 'pdf', '-o', output_file, input_file]
            process = subprocess.Popen(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE
            )
            stdout, stderr = process.communicate()

            return os.path.exists(output_file)
        except Exception as e:
            logger.error(f"Unoconv conversion failed: {str(e)}")
            return False


    def convert_with_pandoc(self, input_file: str, output_file: str) -> bool:
        """
        ä½¿ç”¨Pandocè¿›è¡Œæ–‡ä»¶è½¬æ¢

        Args:
            input_file: è¾“å…¥æ–‡ä»¶è·¯å¾„
            output_file: è¾“å‡ºæ–‡ä»¶è·¯å¾„

        Returns:
            bool: è½¬æ¢æ˜¯å¦æˆåŠŸ
        """
        try:
            cmd = ['pandoc', input_file, '-o', output_file]
            process = subprocess.Popen(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE
            )
            process.communicate()

            return os.path.exists(output_file)
        except Exception as e:
            logger.error(f"Pandoc conversion failed: {str(e)}")
            return False


    @staticmethod
    def get_system_info() -> dict:
        """
        è·å–ç³»ç»Ÿç¯å¢ƒä¿¡æ¯

        Returns:
            dict: ç³»ç»Ÿä¿¡æ¯å­—å…¸ï¼ŒåŒ…å«ï¼š
            - å¹³å°ä¿¡æ¯
            - ç³»ç»Ÿæ¶æ„
            - å¯ç”¨è½¬æ¢å™¨åˆ—è¡¨

        åŠŸèƒ½ï¼š
        1. æ£€æµ‹ç³»ç»Ÿç¯å¢ƒ
        2. æ£€æŸ¥å¯ç”¨è½¬æ¢å™¨
        3. æä¾›ç³»ç»Ÿå…¼å®¹æ€§ä¿¡æ¯
        """
        info = {
            'platform': platform.system(),
            'architecture': platform.architecture(),
            'available_converters': []
        }

        # æ£€æŸ¥LibreOffice
        try:
            if platform.system() == 'Windows':
                libreoffice_paths = [
                    'C:\\Program Files\\LibreOffice\\program\\soffice.exe',
                    'C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe'
                ]
                if any(os.path.exists(path) for path in libreoffice_paths):
                    info['available_converters'].append('LibreOffice')
            else:
                process = subprocess.Popen(['which', 'soffice'], stdout=subprocess.PIPE)
                if process.communicate()[0]:
                    info['available_converters'].append('LibreOffice')
        except:
            pass

        # æ£€æŸ¥å…¶ä»–è½¬æ¢å™¨
        for converter in ['unoconv', 'pandoc']:
            try:
                process = subprocess.Popen(['which', converter], stdout=subprocess.PIPE)
                if process.communicate()[0]:
                    info['available_converters'].append(converter.title())
            except:
                pass

        return info


    def pdf_to_pdfa(self, pdf_path: str) -> Optional[str]:
        """
        å°†PDFè½¬æ¢ä¸ºPDF/Aæ ¼å¼ï¼ˆå­˜æ¡£çº§PDFï¼‰

        Args:
            pdf_path: æºPDFæ–‡ä»¶è·¯å¾„

        Returns:
            str: ç”Ÿæˆçš„PDF/Aæ–‡ä»¶è·¯å¾„ï¼Œå¤±è´¥æ—¶è¿”å›None

        åŠŸèƒ½ç‰¹ç‚¹ï¼š
        1. ç¬¦åˆPDF/A-2æ ‡å‡†
        2. æ”¯æŒé¢œè‰²è½¬æ¢
        3. ç¡®ä¿é•¿æœŸä¿å­˜
        4. ä¿æŒæ–‡æ¡£å®Œæ•´æ€§

        æ³¨æ„äº‹é¡¹ï¼š
        - éœ€è¦å®‰è£…Ghostscript
        - ç¡®ä¿é¢œè‰²é…ç½®æ­£ç¡®
        - éªŒè¯PDF/Aå…¼å®¹æ€§
        """
        try:
            import pypdf
            # ç”Ÿæˆè¾“å‡ºè·¯å¾„
            output_path = self._generate_output_path("converted", "pdf")

            # é…ç½®Ghostscriptå‘½ä»¤è¡Œå‚æ•°
            gs_command = [
                'gs',  # Ghostscriptå‘½ä»¤
                '-dPDFA=2',  # PDF/A-2æ ‡å‡†
                '-dBATCH',  # æ‰¹å¤„ç†æ¨¡å¼
                '-dNOPAUSE',  # ä¸æš‚åœ
                '-sColorConversionStrategy=UseDeviceIndependentColor',  # é¢œè‰²è½¬æ¢ç­–ç•¥
                '-sDEVICE=pdfwrite',  # è¾“å‡ºè®¾å¤‡
                '-dPDFACompatibilityPolicy=1',  # PDF/Aå…¼å®¹æ€§ç­–ç•¥
                f'-sOutputFile={output_path}',  # è¾“å‡ºæ–‡ä»¶
                pdf_path  # è¾“å…¥æ–‡ä»¶
            ]

            # æ‰§è¡Œè½¬æ¢å‘½ä»¤
            result = subprocess.run(gs_command, capture_output=True, text=True)
            if result.returncode != 0:
                raise Exception(f"Ghostscripté”™è¯¯: {result.stderr}")

            outputData(data=f"PDF/Aæ–‡ä»¶å·²ä¿å­˜è‡³: {output_path}", printInfo=self.printInfo)
            return output_path

        except Exception as e:
            logger.error(f"Failed to convert PDF to PDF/A: {str(e)}")
            return None


    async def url_to_pdf(self, url: str) -> Optional[str]:
        """
        å°†ç½‘é¡µè½¬æ¢ä¸ºPDFæ–‡ä»¶ï¼ˆå¼‚æ­¥ç‰ˆæœ¬ï¼‰

        Args:
            url: å¾…è½¬æ¢çš„ç½‘é¡µURL

        Returns:
            str: ç”Ÿæˆçš„PDFæ–‡ä»¶è·¯å¾„ï¼Œå¤±è´¥æ—¶è¿”å›None

        åŠŸèƒ½ï¼š
        1. ä½¿ç”¨å¼‚æ­¥Playwrightè¿›è¡Œç½‘é¡µæ¸²æŸ“
        2. è‡ªåŠ¨å¤„ç†åŠ¨æ€å†…å®¹
        3. æ”¯æŒè‡ªå®šä¹‰PDFå‚æ•°
        4. æä¾›è¯¦ç»†çš„è½¬æ¢çŠ¶æ€
        """
        try:
            from playwright.async_api import async_playwright
            import asyncio

            # ç”Ÿæˆè¾“å‡ºè·¯å¾„
            output_path = self._generate_output_path("url_converted", "pdf")

            # å®‰è£…æµè§ˆå™¨ï¼ˆå¦‚æœéœ€è¦ï¼‰
            try:
                process = await asyncio.create_subprocess_exec(
                    sys.executable, "-m", "playwright", "install", "chromium",
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE
                )
                await process.communicate()
            except Exception as e:
                logger.error(f"Failed to install browsers: {e}")
                return None

            async with async_playwright() as p:
                # å¯åŠ¨æµè§ˆå™¨
                browser = await p.chromium.launch(
                    headless=True,
                    args=['--no-sandbox']
                )

                # åˆ›å»ºä¸Šä¸‹æ–‡å’Œé¡µé¢
                context = await browser.new_context(
                    viewport={'width': 1920, 'height': 1080},
                    device_scale_factor=1.5
                )
                page = await context.new_page()

                # åŠ è½½é¡µé¢
                outputData(f"Loading URL: {url}", self.printInfo)
                await page.goto(url, wait_until="networkidle", timeout=60000)

                # PDFè®¾ç½®
                pdf_options = {
                    "path": output_path,
                    "format": "a4",
                    "margin": {
                        "top": "1cm",
                        "right": "1cm",
                        "bottom": "1cm",
                        "left": "1cm"
                    },
                    "print_background": True,
                    "scale": 0.8
                }

                # ç”ŸæˆPDF
                await page.pdf(**pdf_options)

                # æ¸…ç†èµ„æº
                await context.close()
                await browser.close()

                outputData(f"PDF saved to: {output_path}", self.printInfo)
                return output_path

        except Exception as e:
            logger.error(f"URL to PDF conversion failed: {str(e)}")
            logger.error(f"Detailed error: {traceback.format_exc()}")
            return None


    async def pdf_to_word(self, pdf_path: str, output_format: str = 'docx') -> Optional[str]:
        """
        å°†PDFè½¬æ¢ä¸ºWordæ–‡æ¡£

        Args:
            pdf_path: PDFæ–‡ä»¶è·¯å¾„
            output_format: è¾“å‡ºæ ¼å¼ï¼Œé»˜è®¤ä¸ºdocx

        Returns:
            str: ç”Ÿæˆçš„Wordæ–‡æ¡£è·¯å¾„ï¼Œå¤±è´¥æ—¶è¿”å›None
        """
        try:
            from pdf2docx import Converter
            import asyncio
            import os

            full_path = os.path.join(os.getcwd(), self.upload_dir, pdf_path)
            if os.path.exists(full_path):
                full_path = full_path
            else:
                full_path = os.path.join(os.getcwd(), self.sub_upload_dir, pdf_path)

            # è§„èŒƒåŒ–è·¯å¾„ï¼ˆå¤„ç†æ–œæ /åæ–œæ ï¼‰
            full_file_path = os.path.normpath(full_path)

            # éªŒè¯æ–‡ä»¶æ˜¯å¦å­˜åœ¨
            if not os.path.exists(full_file_path):
                # åˆ—å‡ºuploadç›®å½•ä¸­çš„æ‰€æœ‰æ–‡ä»¶
                upload_dir = os.path.join(os.getcwd(), self.sub_upload_dir)
                if os.path.exists(upload_dir):
                    files = os.listdir(upload_dir)
                    logger.info(f"Uploadç›®å½•ä¸­çš„æ–‡ä»¶: {files}")
                raise FileNotFoundError(f"æ–‡ä»¶æœªæ‰¾åˆ°: {full_file_path}")

            if os.path.getsize(full_file_path) == 0:
                raise ValueError(f"PDFæ–‡ä»¶ä¸ºç©º: {full_file_path}")

            # ç¡®ä¿outputç›®å½•å­˜åœ¨
            output_dir = os.path.join(os.getcwd(), self.output_dir)
            os.makedirs(output_dir, exist_ok=True)

            # ç”Ÿæˆè¾“å‡ºè·¯å¾„
            output_path = self._generate_output_path("converted", output_format)

            # åˆ›å»ºè½¬æ¢å™¨å¯¹è±¡
            cv = Converter(full_file_path)

            try:
                # æ‰§è¡Œè½¬æ¢
                await asyncio.to_thread(cv.convert, output_path)
            finally:
                # ç¡®ä¿åœ¨è½¬æ¢å®Œæˆåå…³é—­è½¬æ¢å™¨
                await asyncio.to_thread(cv.close)

            # éªŒè¯è¾“å‡ºæ–‡ä»¶
            if not os.path.exists(output_path):
                raise RuntimeError("è½¬æ¢å®Œæˆä½†è¾“å‡ºæ–‡ä»¶æœªç”Ÿæˆ")

            if os.path.getsize(output_path) == 0:
                raise RuntimeError("è½¬æ¢å®Œæˆä½†è¾“å‡ºæ–‡ä»¶ä¸ºç©º")

            outputData(f"Wordæ–‡æ¡£å·²ä¿å­˜è‡³: {output_path}", self.printInfo)
            return output_path

        except FileNotFoundError as e:
            logger.error(f"Failed to convert PDF to Word: {str(e)}")
            return None
        except ValueError as e:
            logger.error(f"Failed to convert PDF to Word: {str(e)}")
            return None
        except Exception as e:
            logger.error(f"Failed to convert PDF to Word: {str(e)}")
            import traceback
            logger.error(f"Detailed error: {traceback.format_exc()}")
            return None


    async def pdf_to_text(self, pdf_path: str) -> Optional[str]:
        """
        å°†PDFè½¬æ¢ä¸ºæ–‡æœ¬æ–‡ä»¶

        Args:
            pdf_path: PDFæ–‡ä»¶è·¯å¾„

        Returns:
            str: ç”Ÿæˆçš„æ–‡æœ¬æ–‡ä»¶è·¯å¾„ï¼Œå¤±è´¥æ—¶è¿”å›None

        åŠŸèƒ½ï¼š
        1. æå–æ–‡æœ¬å†…å®¹
        2. ä¿æŒåŸºæœ¬æ’ç‰ˆç»“æ„
        3. æ”¯æŒå¤šè¯­è¨€
        """
        try:
            import pymupdf4llm
            import aiofiles
            import asyncio

            # ç”Ÿæˆè¾“å‡ºè·¯å¾„
            output_path = self._generate_output_path("converted", "txt")

            full_path = os.path.join(os.getcwd(), self.upload_dir, pdf_path)
            if os.path.exists(full_path):
                full_path = os.path.join(self.upload_dir, pdf_path)
            else:
                full_path = os.path.join(self.sub_upload_dir, pdf_path)

            # åœ¨çº¿ç¨‹æ± ä¸­æ‰§è¡ŒPDFè½¬æ¢
            text = await asyncio.to_thread(pymupdf4llm.to_markdown, full_path)

            # å¼‚æ­¥å†™å…¥æ–‡ä»¶
            async with aiofiles.open(output_path, 'w', encoding='utf-8') as f:
                await f.write(text)

            outputData(f"Text files saved to: {output_path}", self.printInfo)
            return output_path

        except Exception as e:
            logger.error(f"PDF to text conversion failed: {str(e)}")
            return None


    def pdf_to_presentation(self, pdf_path: str, output_format: str = 'pptx') -> Optional[str]:
        """
        å°†PDFæ–‡æ¡£è½¬æ¢ä¸ºæ¼”ç¤ºæ–‡ç¨¿ï¼ˆPPTï¼‰

        Args:
            pdf_path: PDFæ–‡ä»¶è·¯å¾„
            output_format: è¾“å‡ºæ ¼å¼ï¼Œé»˜è®¤ä¸º'pptx'

        Returns:
            str: ç”Ÿæˆçš„æ¼”ç¤ºæ–‡ç¨¿è·¯å¾„ï¼Œå¤±è´¥æ—¶è¿”å›None

        åŠŸèƒ½ç‰¹ç‚¹ï¼š
        1. ä¿æŒPDFé¡µé¢å¸ƒå±€
        2. è‡ªåŠ¨å›¾ç‰‡å°ºå¯¸è°ƒæ•´
        3. å±…ä¸­æ˜¾ç¤ºå†…å®¹
        4. é«˜è´¨é‡å›¾åƒè½¬æ¢

        å·¥ä½œæµç¨‹ï¼š
        1. PDFè½¬æ¢ä¸ºé«˜è´¨é‡å›¾ç‰‡
        2. åˆ›å»ºPPTå¹»ç¯ç‰‡
        3. å°†å›¾ç‰‡æ·»åŠ åˆ°å¹»ç¯ç‰‡
        4. ä¼˜åŒ–å¸ƒå±€å’Œæ˜¾ç¤º
        """
        try:
            # è·å–Popplerå·¥å…·è·¯å¾„
            poppler_path = os.path.join(self.base_dir, 'poppler', 'bin')

            # ç”Ÿæˆè¾“å‡ºæ–‡ä»¶è·¯å¾„
            output_path = self._generate_output_path("converted", output_format)

            # åˆå§‹åŒ–PPTæ–‡æ¡£
            prs = Presentation()

            full_path = os.path.join(os.getcwd(), self.upload_dir, pdf_path)
            if os.path.exists(full_path):
                full_path = os.path.join(self.upload_dir, pdf_path)
            else:
                full_path = os.path.join(self.sub_upload_dir, pdf_path)

            # å°†PDFé¡µé¢è½¬æ¢ä¸ºé«˜è´¨é‡å›¾ç‰‡ï¼ˆ300dpiï¼‰
            images = convert_from_path(
                full_path,
                dpi=300,
                poppler_path=poppler_path
            )

            # è·å–PPTé¡µé¢å°ºå¯¸
            slide_width = prs.slide_width  # å¹»ç¯ç‰‡å®½åº¦
            slide_height = prs.slide_height  # å¹»ç¯ç‰‡é«˜åº¦

            # å¤„ç†æ¯ä¸€é¡µ
            for image in images:
                # åˆ›å»ºç©ºç™½å¹»ç¯ç‰‡ï¼ˆä½¿ç”¨å¸ƒå±€6ï¼šç©ºç™½å¸ƒå±€ï¼‰
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)

                # åˆ›å»ºä¸´æ—¶å›¾ç‰‡æ–‡ä»¶
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                    tmp_name = tmp.name  # ä¿å­˜ä¸´æ—¶æ–‡ä»¶å
                    tmp.close()  # å…³é—­æ–‡ä»¶å¥æŸ„

                    # å°†å›¾ç‰‡ä¿å­˜ä¸ºPNGæ ¼å¼
                    image.save(tmp_name, 'PNG')

                    # è®¾ç½®å›¾ç‰‡åœ¨å¹»ç¯ç‰‡ä¸­çš„å°ºå¯¸
                    img_width = Inches(10)  # æ ‡å‡†å®½åº¦10è‹±å¯¸
                    img_height = Inches(7.5)  # æ ‡å‡†é«˜åº¦7.5è‹±å¯¸

                    # è®¡ç®—å±…ä¸­ä½ç½®
                    left = (slide_width - img_width) / 2
                    top = (slide_height - img_height) / 2

                    # å°†å›¾ç‰‡æ·»åŠ åˆ°å¹»ç¯ç‰‡
                    slide.shapes.add_picture(
                        tmp_name,  # å›¾ç‰‡è·¯å¾„
                        left=left,  # å·¦è¾¹è·
                        top=top,  # ä¸Šè¾¹è·
                        width=img_width,  # å›¾ç‰‡å®½åº¦
                        height=img_height  # å›¾ç‰‡é«˜åº¦
                    )

                    # æ¸…ç†ä¸´æ—¶æ–‡ä»¶
                    os.unlink(tmp_name)

            # ä¿å­˜æœ€ç»ˆçš„PPTæ–‡ä»¶
            prs.save(output_path)
            outputData(data=f"æ¼”ç¤ºæ–‡ç¨¿å·²ä¿å­˜è‡³: {output_path}", printInfo=self.printInfo)
            return output_path

        except Exception as e:
            logger.error(f"Failed to convert PDF to presentation: {str(e)}")
            return None


    def pdf_to_image(self, pdf_path: str, image_format: str = 'png',
                     single_or_multiple: str = 'multiple',
                     color_type: str = 'color', dpi: str = '300') -> Optional[str]:
        """
        å°†PDFè½¬æ¢ä¸ºå›¾ç‰‡

        Args:
            pdf_path: PDFæ–‡ä»¶è·¯å¾„
            image_format: å›¾ç‰‡æ ¼å¼ï¼ˆpng/jpgï¼‰
            single_or_multiple: å•é¡µæˆ–å¤šé¡µæ¨¡å¼
            color_type: é¢œè‰²æ¨¡å¼
            dpi: å›¾åƒåˆ†è¾¨ç‡

        Returns:
            str: ç”Ÿæˆçš„å›¾ç‰‡æ–‡ä»¶æˆ–ZIPæ–‡ä»¶è·¯å¾„
        """
        try:
            if not self._check_poppler():
                logger.info("\nPoppler is not properly installed.")
                self._download_poppler()
                if not self._check_poppler():
                    return None

            from pdf2image import convert_from_path
            import tempfile
            import shutil

            # ç¡®ä¿Poppleråœ¨ç¯å¢ƒå˜é‡ä¸­
            os.environ['PATH'] = f"{self.poppler_path};{os.environ['PATH']}"

            # æ„å»ºå®Œæ•´çš„PDFè·¯å¾„
            full_path = os.path.join(os.getcwd(), self.upload_dir, pdf_path)
            if os.path.exists(full_path):
                full_path = os.path.join(self.upload_dir, pdf_path)
            else:
                full_path = os.path.join(self.sub_upload_dir, pdf_path)
            # è§„èŒƒåŒ–è·¯å¾„ï¼ˆå¤„ç†æ–œæ /åæ–œæ ï¼‰
            full_pdf_path = os.path.normpath(full_path)

            # æ£€æŸ¥æ–‡ä»¶æ˜¯å¦å­˜åœ¨
            if not os.path.exists(full_pdf_path):
                logger.info(f"æ–‡ä»¶ä¸å­˜åœ¨: {full_pdf_path}")
                return None

            # åˆ›å»ºä¸´æ—¶ç›®å½•
            temp_dir = tempfile.mkdtemp()
            try:
                # è½¬æ¢PDFé¡µé¢
                images = convert_from_path(
                    full_pdf_path,
                    dpi=int(dpi),
                    fmt=image_format.lower(),
                    output_folder=temp_dir,
                    poppler_path=self.poppler_path
                )

                if single_or_multiple == 'multiple':
                    # åˆ›å»ºZIPåŒ…å«æ‰€æœ‰é¡µé¢
                    output_path = self._generate_output_path("converted", "zip")
                    with zipfile.ZipFile(output_path, 'w') as zf:
                        for i, image in enumerate(images):
                            image_path = os.path.join(temp_dir, f'page_{i + 1}.{image_format}')
                            image.save(image_path, format=image_format.upper())
                            zf.write(image_path, f'page_{i + 1}.{image_format}')
                else:
                    # ä»…ä¿å­˜ç¬¬ä¸€é¡µ
                    output_path = self._generate_output_path("converted", image_format)
                    images[0].save(output_path, format=image_format.upper())

                outputData(f"Images saved to: {output_path}", self.printInfo)
                return output_path

            finally:
                # æ¸…ç†ä¸´æ—¶æ–‡ä»¶
                try:
                    shutil.rmtree(temp_dir)
                except Exception as e:
                    logger.info(f"Warning: Failed to clean up temporary directory: {e}")

        except Exception as e:
            logger.error(f"\nPDF to image conversion failed: {str(e)}")
            logger.error(f"Detailed error: {traceback.format_exc()}")
            return None


    def pdf_to_markdown(self, pdf_path: str, output_dir: str = "output"):
        """
        å°†PDFè½¬æ¢ä¸ºMarkdownæ ¼å¼

        Args:
            pdf_path: PDFæ–‡ä»¶è·¯å¾„
            output_dir: è¾“å‡ºç›®å½•

        Returns:
            str: ç”Ÿæˆçš„Markdownæ–‡ä»¶è·¯å¾„

        åŠŸèƒ½ï¼š
        1. ä¿æŒæ–‡æ¡£ç»“æ„
        2. è½¬æ¢æ ‡é¢˜å’Œåˆ—è¡¨
        3. å¤„ç†å›¾ç‰‡å’Œé“¾æ¥
        4. æ”¯æŒä»£ç å—å’Œè¡¨æ ¼
        """
        try:
            import pymupdf4llm
            from pathlib import Path

            full_path = os.path.join(os.getcwd(), self.upload_dir, pdf_path)
            if os.path.exists(full_path):
                full_path = os.path.join(self.upload_dir, pdf_path)
            else:
                full_path = os.path.join(self.sub_upload_dir, pdf_path)
            # è§„èŒƒåŒ–è·¯å¾„ï¼ˆå¤„ç†æ–œæ /åæ–œæ ï¼‰
            full_pdf_path = os.path.normpath(full_path)
            # è½¬æ¢å†…å®¹
            text = pymupdf4llm.to_markdown(full_pdf_path)

            output_path = self._generate_output_path("converted", "md")
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)

            outputData(f"Markdown files saved to: {output_path}", self.printInfo)
            return output_path

        except Exception as e:
            logger.error(f"PDF to Markdown conversion failed: {str(e)}")
            return None


    def markdown_to_pdf(self, markdown_path: str) -> Optional[str]:
        """
        å°†Markdownè½¬æ¢ä¸ºPDF

        Args:
            markdown_path: Markdownæ–‡ä»¶è·¯å¾„

        Returns:
            str: ç”Ÿæˆçš„PDFæ–‡ä»¶è·¯å¾„

        åŠŸèƒ½ï¼š
        1. æ”¯æŒå®Œæ•´çš„Markdownè¯­æ³•
        2. è‡ªå®šä¹‰æ ·å¼å’Œä¸»é¢˜
        3. ä»£ç é«˜äº®
        4. è¡¨æ ¼å’Œå›¾ç‰‡å¤„ç†
        """
        try:
            import markdown
            import pdfkit

            # é…ç½®wkhtmltopdfè·¯å¾„
            path_wkhtmltopdf = r"C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe"
            if not os.path.isfile(path_wkhtmltopdf):
                logger.info(f"Error: wkhtmltopdf not found at '{path_wkhtmltopdf}'")
                logger.info("Please install from https://wkhtmltopdf.org/downloads.html")
                return None

            config = pdfkit.configuration(wkhtmltopdf=path_wkhtmltopdf)
            output_path = self._generate_output_path("converted", "pdf")
            full_path = os.path.join(os.getcwd(), self.upload_dir, markdown_path)
            if os.path.exists(full_path):
                full_path = os.path.join(self.upload_dir, markdown_path)
            else:
                full_path = os.path.join(self.sub_upload_dir, markdown_path)
            # è§„èŒƒåŒ–è·¯å¾„ï¼ˆå¤„ç†æ–œæ /åæ–œæ ï¼‰
            full_md_path = os.path.normpath(full_path)
            # è¯»å–å’Œè½¬æ¢Markdown
            with open(full_md_path, 'r', encoding='utf-8') as f:
                md_content = f.read()

            # è½¬æ¢ä¸ºHTML
            html_content = markdown.markdown(
                md_content,
                extensions=['tables', 'fenced_code', 'codehilite']
            )

            # æ·»åŠ æ ·å¼
            html_doc = f'''
                <!DOCTYPE html>
                <html>
                <head>
                    <meta charset="utf-8">
                    <style>
                        body {{ font-family: Arial, sans-serif; margin: 2cm; }}
                        h1, h2, h3 {{ color: #333; }}
                        code {{ background-color: #f5f5f5; padding: 2px 4px; }}
                        pre {{ background-color: #f5f5f5; padding: 1em; }}
                        blockquote {{ border-left: 4px solid #ccc; margin-left: 0; padding-left: 1em; }}
                    </style>
                </head>
                <body>{html_content}</body>
                </html>
            '''

            # ä¿å­˜ä¸´æ—¶HTML
            temp_html = os.path.join(self.output_directory, 'temp.html')
            with open(temp_html, 'w', encoding='utf-8') as f:
                f.write(html_doc)

            # é…ç½®PDFç”Ÿæˆé€‰é¡¹
            options = {
                'enable-local-files-access': True,
                'no-stop-slow-scripts': True,
            }

            # è½¬æ¢ä¸ºPDF
            pdfkit.from_file(temp_html, output_path, configuration=config, options=options)

            # æ¸…ç†ä¸´æ—¶æ–‡ä»¶
            os.remove(temp_html)

            outputData(f"PDF saved to: {output_path}", self.printInfo)
            return output_path

        except Exception as e:
            logger.error(f"Markdown to PDF conversion failed: {str(e)}")
            return None


    def file_to_pdf(self, file_path: str) -> Optional[str]:
        """
        é€šç”¨æ–‡ä»¶è½¬PDFæ–¹æ³•

        Args:
            file_path: æºæ–‡ä»¶è·¯å¾„

        Returns:
            str: ç”Ÿæˆçš„PDFæ–‡ä»¶è·¯å¾„

        åŠŸèƒ½ï¼š
        1. æ”¯æŒå¤šç§æ–‡ä»¶æ ¼å¼
        2. å¤šå¼•æ“è‡ªåŠ¨åˆ‡æ¢
        3. æ™ºèƒ½å¤±è´¥æ¢å¤
        4. è¯¦ç»†çš„è½¬æ¢æ—¥å¿—
        """
        try:
            full_path = os.path.join(os.getcwd(), self.upload_dir, file_path)
            if os.path.exists(full_path):
                full_path = os.path.join(self.upload_dir, file_path)
            else:
                full_path = os.path.join(self.sub_upload_dir, file_path)
            # è§„èŒƒåŒ–è·¯å¾„ï¼ˆå¤„ç†æ–œæ /åæ–œæ ï¼‰
            full_file_path = os.path.normpath(full_path)
            output_path = self._generate_output_path("converted", "pdf")

            # æ£€æŸ¥æ–‡ä»¶å­˜åœ¨æ€§
            if not os.path.exists(full_file_path):
                logger.info(f"File not found: {full_file_path}")
                return None

            # å°è¯•ä¸åŒçš„è½¬æ¢æ–¹æ³•
            converters = [
                (self.convert_with_libreoffice, "LibreOffice"),
                (self.convert_with_unoconv, "Unoconv"),
                (self.convert_with_pandoc, "Pandoc")
            ]

            # ä¾æ¬¡å°è¯•æ¯ä¸ªè½¬æ¢å™¨
            for converter, name in converters:
                try:
                    if converter(full_file_path, output_path):
                        outputData(f"PDF saved to: {output_path}", self.printInfo)
                        return output_path
                except Exception:
                    continue

            logger.info("File to PDF conversion failed")
            return None

        except Exception as e:
            logger.error(f"File to PDF conversion failed: {str(e)}")
            logger.error(f"Detailed error: {traceback.format_exc()}")
            return None


    async def run(self, **kwargs) -> str:
        """
        æ‰§è¡Œæ–‡ä»¶è½¬æ¢çš„ç»Ÿä¸€å…¥å£

        Args:
            **kwargs: å…³é”®å­—å‚æ•°
                - conversion_type: è½¬æ¢ç±»å‹
                - input_path/url: è¾“å…¥æ–‡ä»¶è·¯å¾„æˆ–URL

        Returns:
            str: è½¬æ¢ç»“æœæˆ–é”™è¯¯ä¿¡æ¯
        """
        try:
            # è·å–å‚æ•°
            conversion_type = kwargs.get("conversion_type")
            input_path = kwargs.get("input_path") or kwargs.get("url")  # æ”¯æŒ url å‚æ•°

            # å‚æ•°éªŒè¯
            if not conversion_type:
                return "è½¬æ¢å¤±è´¥: ç¼ºå°‘è½¬æ¢ç±»å‹å‚æ•°"
            if not input_path:
                return "è½¬æ¢å¤±è´¥: ç¼ºå°‘è¾“å…¥è·¯å¾„å‚æ•°"

            # æ ¹æ®ç±»å‹è°ƒç”¨ç›¸åº”æ–¹æ³•
            conversion_methods = {
                "url_to_pdf": self.url_to_pdf,
                "pdf_to_word": self.pdf_to_word,
                "pdf_to_text": self.pdf_to_text,
                "pdf_to_image": lambda x: self.pdf_to_image(x, single_or_multiple='multiple'),
                "pdf_to_ppt": self.pdf_to_presentation,
                "pdf_to_markdown": self.pdf_to_markdown,
                "markdown_to_pdf": self.markdown_to_pdf,
                "file_to_pdf": self.file_to_pdf
            }

            converter = conversion_methods.get(conversion_type)
            if not converter:
                return f"ä¸æ”¯æŒçš„è½¬æ¢ç±»å‹: {conversion_type}"

            # æ ¹æ®æ–¹æ³•æ˜¯å¦æ˜¯åç¨‹å†³å®šè°ƒç”¨æ–¹å¼
            if asyncio.iscoroutinefunction(converter):
                result = await converter(input_path)
            else:
                result = converter(input_path)

            # å¤„ç†è½¬æ¢ç»“æœ
            if result:
                return result
            else:
                return "ğŸ“„ è½¬æ¢çŠ¶æ€ï¼š\nè½¬æ¢å¤±è´¥\nè¾“å‡ºè·¯å¾„ï¼šæœªç”Ÿæˆ"

        except Exception as e:
            return f"è½¬æ¢å¤±è´¥: {str(e)}"
